#!/usr/bin/env python3
"""Gera templates padrão do HermesOffice a partir dos estilos da ppt-agent.

3 looks: nocturne (dark violet Linear), blue-white (executivo Apple),
mocha (editorial Anthropic). Mesma convenção de elementos do deck-base:
  Slide 1 capa: "Title 1" + "Subtitle 2"
  Slide 2 conteúdo: "Title 1" + "Content Placeholder 2"
  Slide 3 seção: "Title 1" + "Text Placeholder 2"
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = "/tmp/genoffice/templates"

LOOKS = [
    dict(name="deck-nocturne.pptx", bg=RGBColor(0x0A, 0x06, 0x12), accent=RGBColor(0x9B, 0x64, 0xFF),
         text=RGBColor(0xFF, 0xFF, 0xFF), sub=RGBColor(0xC4, 0xA8, 0xFF),
         title="Nocturne", tag="Dark · Linear"),
    dict(name="deck-blue-white.pptx", bg=RGBColor(0xFF, 0xFF, 0xFF), accent=RGBColor(0x25, 0x63, 0xEB),
         text=RGBColor(0x0A, 0x1D, 0x3A), sub=RGBColor(0x64, 0x74, 0x8B),
         title="Blue White", tag="Clean · Executive"),
    dict(name="deck-mocha.pptx", bg=RGBColor(0xF3, 0xEB, 0xDE), accent=RGBColor(0xC1, 0x50, 0x2E),
         text=RGBColor(0x2A, 0x18, 0x10), sub=RGBColor(0x8B, 0x5A, 0x3C),
         title="Mocha", tag="Editorial · Anthropic"),
]


def style_run(run, color, size, bold=False, italic=False):
    f = run.font
    f.color.rgb = color
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = "Inter"


def add_accent_bar(slide, accent, y=Inches(7.35)):
    from pptx.util import Inches as I
    bar = slide.shapes.add_shape(1, 0, y, I(13.333), I(0.08))  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False


def build(look):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: capa ---
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = look["bg"]
    s.shapes.title.text = look["title"]
    style_run(s.shapes.title.text_frame.paragraphs[0].runs[0], look["accent"], 60, bold=True)
    s.placeholders[1].text = f"Subtítulo — template {look['tag']} · HermesOffice"
    style_run(s.placeholders[1].text_frame.paragraphs[0].runs[0], look["sub"], 24)
    add_accent_bar(s, look["accent"])

    # --- Slide 2: título + conteúdo ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = look["bg"]
    s.shapes.title.text = "Título da seção"
    style_run(s.shapes.title.text_frame.paragraphs[0].runs[0], look["text"], 36, bold=True)
    tf = s.placeholders[1].text_frame
    for i, line in enumerate(["Primeiro ponto — destaque do argumento",
                              "Segundo ponto — dado ou evidência",
                              "Terceiro ponto — conclusão parcial"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        style_run(run, look["sub"] if i else look["text"], 20, bold=(i == 0))
    add_accent_bar(s, look["accent"])

    # --- Slide 3: seção ---
    s = prs.slides.add_slide(prs.slide_layouts[2])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = look["bg"]
    s.shapes.title.text = "Seção de destaque"
    style_run(s.shapes.title.text_frame.paragraphs[0].runs[0], look["accent"], 44, bold=True)
    if s.placeholders and len(s.placeholders) > 1:
        ph = s.placeholders[1]
        ph.text = "Capítulo"
        style_run(ph.text_frame.paragraphs[0].runs[0], look["sub"], 18)
    add_accent_bar(s, look["accent"])

    path = f"{OUT}/{look['name']}"
    prs.save(path)
    print(f"OK  {path}")


for look in LOOKS:
    build(look)
print("templates gerados")
